#!/usr/bin/env python3
"""Generate a PowerPoint (.pptx) file from a JSON intermediate format.

This script is the PPTX generation backend for the Presentation Composer skill.
It reads a structured JSON file produced by the Composer (Step 4) and generates
a .pptx file with proper slide layout mapping, template support, and branding.

Usage:
    python generate_pptx.py --input composed-draft.json --output presentation.pptx
    python generate_pptx.py --input draft.json --output deck.pptx --template brand.pptx
    python generate_pptx.py --input draft.json --output deck.pptx --font Arial --accent-color "#ff6600"

Requirements:
    pip install python-pptx

JSON Input Format (per architecture ADR-01):
    {
      "slides": [
        {
          "number": 1,
          "title": "Slide Title",
          "layout": "title|content|metrics|comparison|cta|timeline|architecture",
          "body": ["bullet 1", "bullet 2"],
          "table": null | {"headers": [...], "rows": [[...]]},
          "speaker_notes": null | "Notes text",
          "citations": ["artifact.md"],
          "mermaid": null | "graph TD; ..."
        }
      ],
      "metadata": {
        "type": "sprint-review",
        "date": "2026-04-04",
        "project": "Project Name",
        "audience": "technical",
        "format": "pptx"
      }
    }
"""

# FR-09: Import guard -- check for python-pptx before anything else.
try:
    import pptx
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    import sys

    print(
        "Error: python-pptx is required. Install with: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)

import argparse
import json
import sys
from pathlib import Path
from typing import Any


# --- Constants ---

DEFAULT_FONT = "Calibri"
DEFAULT_ACCENT_COLOR = "#2d5aa0"

# FR-08: Layout name-first, index-fallback mapping (per architecture Section 1.5).
# All non-title layouts use "Title and Content" as the base. Differentiation
# happens in content population, not layout selection.
LAYOUT_MAP: dict[str, dict[str, Any]] = {
    "title": {"name": "Title Slide", "index": 0},
    "content": {"name": "Title and Content", "index": 1},
    "metrics": {"name": "Title and Content", "index": 1},
    "comparison": {"name": "Title and Content", "index": 1},
    "cta": {"name": "Title and Content", "index": 1},
    "timeline": {"name": "Title and Content", "index": 1},
    "architecture": {"name": "Title and Content", "index": 1},
}

# Fallback for unknown layout values.
DEFAULT_LAYOUT = {"name": "Title and Content", "index": 1}


# --- Layout Resolution ---


def resolve_layout(prs: pptx.presentation.Presentation, layout_key: str):
    """Resolve a JSON layout key to a PowerPoint slide layout.

    Strategy (FR-08): name-first, index-fallback.
    1. Try to find a layout by name in the presentation's slide layouts.
    2. If not found, fall back to the mapped index.
    3. If index is also out of range, use layout index 0.
    """
    mapping = LAYOUT_MAP.get(layout_key, DEFAULT_LAYOUT)
    target_name = mapping["name"]
    target_index = mapping["index"]

    # Name-first: search all slide layouts in the first slide master.
    for layout in prs.slide_layouts:
        if layout.name == target_name:
            return layout

    # Index-fallback: use the mapped index.
    try:
        return prs.slide_layouts[target_index]
    except IndexError:
        # Last resort: use whatever layout 0 is.
        return prs.slide_layouts[0]


# --- Color Parsing ---


def parse_hex_color(hex_str: str) -> RGBColor:
    """Parse a hex color string (#RRGGBB or RRGGBB) to an RGBColor."""
    hex_str = hex_str.lstrip("#")
    if len(hex_str) != 6:
        raise ValueError(f"Invalid hex color: #{hex_str}. Expected 6 hex digits.")
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


# --- Slide Population ---


def set_text_frame_font(text_frame, font_name: str, accent_color: RGBColor) -> None:
    """Apply font and color settings to all paragraphs in a text frame."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.color.rgb = accent_color


def populate_title_slide(
    slide, slide_data: dict, font_name: str, accent_color: RGBColor
) -> None:
    """Populate a title slide (layout index 0).

    Title slides typically have a title placeholder and a subtitle placeholder.
    """
    # Set title.
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get("title", "")
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.name = font_name

    # Set subtitle from body (join bullets) or first body item.
    body_items = slide_data.get("body", [])
    subtitle_text = "\n".join(body_items) if body_items else ""

    # Find subtitle placeholder (usually index 1).
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = subtitle_text
            for run in shape.text_frame.paragraphs[0].runs:
                run.font.name = font_name
            break


def add_bullets(text_frame, items: list[str], font_name: str, font_size: int = 18) -> None:
    """Add bullet points to a text frame, one paragraph per item."""
    # Clear default empty paragraph content but keep it for the first bullet.
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.name = font_name


def add_table_to_slide(
    slide, table_data: dict, font_name: str, top: float = 2.5
) -> None:
    """Add a table shape to a slide from structured table data.

    table_data: {"headers": [...], "rows": [[...]]}
    """
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers and not rows:
        return

    col_count = len(headers) if headers else (len(rows[0]) if rows else 0)
    row_count = (1 if headers else 0) + len(rows)

    if col_count == 0 or row_count == 0:
        return

    # Position and size the table.
    left = Inches(0.5)
    top_pos = Inches(top)
    width = Inches(9.0)
    height = Inches(0.4 * row_count)

    table_shape = slide.shapes.add_table(row_count, col_count, left, top_pos, width, height)
    table = table_shape.table

    # Populate header row.
    row_idx = 0
    if headers:
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(header_text)
            for p in cell.text_frame.paragraphs:
                p.font.name = font_name
                p.font.size = Pt(14)
                p.font.bold = True
        row_idx = 1

    # Populate data rows.
    for row_data in rows:
        for col_idx, cell_text in enumerate(row_data):
            if col_idx < col_count:
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_text)
                for p in cell.text_frame.paragraphs:
                    p.font.name = font_name
                    p.font.size = Pt(12)
        row_idx += 1


def populate_content_slide(
    slide,
    slide_data: dict,
    font_name: str,
    accent_color: RGBColor,
    layout_key: str,
) -> None:
    """Populate a content-type slide (Title and Content layout).

    Handles layout-specific content formatting:
    - content: standard bullets
    - metrics: headline finding as title emphasis, data points as body
    - comparison: renders table data
    - cta: numbered list with action items
    - timeline: renders as table with status indicators
    - architecture: bullets + Mermaid fallback note
    """
    # Set title.
    if slide.shapes.title:
        slide.shapes.title.text = slide_data.get("title", "")
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.name = font_name

    body_items = slide_data.get("body", [])
    table_data = slide_data.get("table")
    mermaid = slide_data.get("mermaid")

    # Find the body placeholder (usually index 1).
    body_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            body_placeholder = shape
            break

    if layout_key == "comparison" and table_data:
        # Render comparison as a table shape.
        add_table_to_slide(slide, table_data, font_name)
        # Add summary bullets below the table if body items exist.
        if body_items and body_placeholder:
            add_bullets(body_placeholder.text_frame, body_items, font_name)

    elif layout_key == "timeline" and table_data:
        # Render timeline as a table.
        add_table_to_slide(slide, table_data, font_name)
        if body_items and body_placeholder:
            add_bullets(body_placeholder.text_frame, body_items, font_name)

    elif layout_key == "cta":
        # Numbered list for call-to-action items.
        if body_placeholder and body_items:
            numbered = [f"{i + 1}. {item}" for i, item in enumerate(body_items)]
            add_bullets(body_placeholder.text_frame, numbered, font_name)

    elif layout_key == "architecture":
        # Standard bullets + Mermaid fallback note.
        items = list(body_items)
        if mermaid:
            items.append("[Mermaid diagram -- render separately or paste as image]")
        if body_placeholder and items:
            add_bullets(body_placeholder.text_frame, items, font_name)

    elif layout_key == "metrics":
        # Metrics: emphasize headline data.
        if body_placeholder and body_items:
            add_bullets(body_placeholder.text_frame, body_items, font_name, font_size=20)

    else:
        # Standard content: bullets.
        if body_placeholder and body_items:
            add_bullets(body_placeholder.text_frame, body_items, font_name)

    # If there is table data but no special layout handling, add it below.
    if table_data and layout_key not in ("comparison", "timeline"):
        add_table_to_slide(slide, table_data, font_name, top=4.0)


def add_speaker_notes(slide, notes_text: str | None) -> None:
    """Add speaker notes to a slide's notes pane (OQ-5 resolution)."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


# --- Main Generation ---


def generate_pptx(
    input_path: str,
    output_path: str,
    template_path: str | None = None,
    font_name: str = DEFAULT_FONT,
    accent_color_hex: str = DEFAULT_ACCENT_COLOR,
) -> int:
    """Generate a .pptx file from a JSON intermediate.

    Args:
        input_path: Path to the composed-draft.json file.
        output_path: Path for the generated .pptx file.
        template_path: Optional path to a .pptx template file.
        font_name: Font family name (default: Calibri).
        accent_color_hex: Hex color for accent elements (default: #2d5aa0).

    Returns:
        Number of slides generated.

    Raises:
        FileNotFoundError: If input file or template file does not exist.
        json.JSONDecodeError: If input file is not valid JSON.
        ValueError: If JSON structure is invalid.
    """
    # Parse accent color.
    try:
        accent_color = parse_hex_color(accent_color_hex)
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)

    # Read JSON input.
    input_file = Path(input_path)
    if not input_file.exists():
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    try:
        with open(input_file, "r", encoding="utf-8") as f:
            data = json.load(f)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON in {input_path}: {e}", file=sys.stderr)
        sys.exit(1)

    # Validate structure.
    slides_data = data.get("slides")
    if not slides_data or not isinstance(slides_data, list):
        print(
            'Error: JSON must contain a non-empty "slides" array.',
            file=sys.stderr,
        )
        sys.exit(1)

    # Create presentation -- from template or blank.
    if template_path:
        template_file = Path(template_path)
        if not template_file.exists():
            print(
                f"Error: Template file not found: {template_path}",
                file=sys.stderr,
            )
            sys.exit(1)
        prs = pptx.Presentation(template_path)
    else:
        prs = pptx.Presentation()

    # Set slide dimensions to widescreen 16:9 (only for blank presentations).
    if not template_path:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    # Generate slides.
    slide_count = 0
    for slide_data in slides_data:
        layout_key = slide_data.get("layout", "content").lower()
        layout = resolve_layout(prs, layout_key)
        slide = prs.slides.add_slide(layout)

        if layout_key == "title":
            populate_title_slide(slide, slide_data, font_name, accent_color)
        else:
            populate_content_slide(
                slide, slide_data, font_name, accent_color, layout_key
            )

        # Speaker notes (OQ-5).
        add_speaker_notes(slide, slide_data.get("speaker_notes"))

        slide_count += 1

    # Ensure output directory exists.
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)

    # Save.
    prs.save(output_path)

    return slide_count


# --- CLI ---


def parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description="Generate a .pptx file from a JSON intermediate format.",
        epilog=(
            "Example:\n"
            "  python generate_pptx.py --input composed-draft.json "
            "--output sprint-review-2026-04-04.pptx\n"
            "  python generate_pptx.py --input draft.json --output deck.pptx "
            "--template brand.pptx --font Arial --accent-color '#ff6600'"
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "--input",
        required=True,
        help="Path to the composed-draft.json input file.",
    )
    parser.add_argument(
        "--output",
        required=True,
        help="Path for the generated .pptx output file.",
    )
    parser.add_argument(
        "--template",
        default=None,
        help="Path to a .pptx template file for branding (slide masters, fonts, colors).",
    )
    parser.add_argument(
        "--font",
        default=DEFAULT_FONT,
        help=f"Font family name (default: {DEFAULT_FONT}).",
    )
    parser.add_argument(
        "--accent-color",
        default=DEFAULT_ACCENT_COLOR,
        help=f"Hex accent color, e.g., '#ff6600' (default: {DEFAULT_ACCENT_COLOR}).",
    )
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> None:
    """Entry point for CLI invocation."""
    args = parse_args(argv)

    slide_count = generate_pptx(
        input_path=args.input,
        output_path=args.output,
        template_path=args.template,
        font_name=args.font,
        accent_color_hex=args.accent_color,
    )

    print(f"Saved: {args.output} ({slide_count} slides)")


if __name__ == "__main__":
    main()
